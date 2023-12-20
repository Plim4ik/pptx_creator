from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from texts import slides_content

# Create a presentation object
prs = Presentation()
     
def add_slide(prs, title, content):
    # Add a slide with a title and content layout
    slide_layout = prs.slide_layouts[1]  # Using layout 1 for title and content
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    # Set content
    content_placeholder = slide.placeholders[1]
    content_placeholder.text = content

    # Format title
    for paragraph in title_placeholder.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(24)
            run.font.name = 'Times New Roman'

    # Format content
    for paragraph in content_placeholder.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = Pt(24)
            run.font.name = 'Times New Roman'

# Add slides to the presentation
for title, content in slides_content:
    add_slide(prs, title, content)

# Save the presentation
prs.save('presentation.pptx')